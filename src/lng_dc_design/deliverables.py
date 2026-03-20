from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .thermo import ensure_directory


def _summary_map(output_dir: Path) -> dict[str, dict[str, str]]:
    frame = pd.read_csv(output_dir / "summary.csv")
    return {
        str(row["metric"]): {
            "value": row["value"],
            "unit": row["unit"],
            "source_ids": row["source_ids"],
        }
        for _, row in frame.iterrows()
    }


def _format_number(value: float, digits: int = 1) -> str:
    if pd.isna(value):
        return "-"
    return f"{float(value):,.{digits}f}"


def _parse_sources(sources_path: Path) -> dict[str, dict[str, str]]:
    entries: dict[str, dict[str, str]] = {}
    for line in sources_path.read_text(encoding="utf-8").splitlines():
        if not line.startswith("| SRC-"):
            continue
        parts = [part.strip() for part in line.strip().split("|")[1:-1]]
        if len(parts) < 6:
            continue
        entries[parts[0]] = {
            "title": parts[1],
            "link_or_path": parts[2],
            "type": parts[3],
            "used_values": parts[4],
            "used_in": parts[5],
        }
    return entries


def _markdown_table(frame: pd.DataFrame, columns: list[str]) -> str:
    header = "| " + " | ".join(columns) + " |"
    divider = "| " + " | ".join(["---"] * len(columns)) + " |"
    rows = [header, divider]
    for _, row in frame[columns].iterrows():
        rows.append("| " + " | ".join(str(row[column]) for column in columns) + " |")
    return "\n".join(rows)


def build_report(project_root: Path) -> Path:
    output_dir = project_root / "output"
    deliverables_dir = ensure_directory(project_root / "deliverables")
    summary = _summary_map(output_dir)
    alternatives = pd.read_csv(output_dir / "alternative_designs.csv")
    distance = pd.read_csv(output_dir / "distance_scenarios.csv")
    temperature = pd.read_csv(output_dir / "supply_temperature_sweep.csv")
    annual = pd.read_csv(output_dir / "annual_summary.csv")
    payback = pd.read_csv(output_dir / "payback_allowable_capex.csv")
    legacy = pd.read_csv(output_dir / "legacy_comparison.csv")
    sources = _parse_sources(project_root / "docs" / "sources.md")

    selected_alternative = alternatives.iloc[0]
    long_distance = distance.sort_values("distance_km").iloc[-1]
    best_temp = temperature[temperature["status"] == "feasible"].sort_values("pump_power_kw").iloc[0]
    annual_report = annual.copy()
    annual_report["value"] = annual_report.apply(
        lambda row: _format_number(float(row["value"]) / 1_000_000.0)
        if row["metric"] == "Electricity cost saving"
        else _format_number(float(row["value"])),
        axis=1,
    )

    report_path = deliverables_dir / "report_draft.md"
    report_lines = [
        "# LNG Cold-Energy Based IDC Cooling System Design Study",
        "",
        "## Abstract",
        "",
        (
            "This report reconstructs the 2022 thermal system design project as a reproducible Python-based study for a "
            "large-scale internet data center cooled by LNG cold energy. The final design delivers a modeled cooling load of "
            f"**{_format_number(float(summary['IDC total cooling load']['value']))} kW** while reducing the reference "
            f"R-134a compressor power requirement from **{_format_number(float(summary['Baseline R-134a compressor power']['value']))} kW** "
            f"to an LNG loop pumping demand of **{_format_number(float(summary['LNG system pump power']['value']))} kW**. "
            f"The selected coolant is **{summary['Selected coolant']['value']}**, the LNG vaporizer pinch constraint is "
            f"held at **10.0 K**, and the estimated annual electricity saving is "
            f"**{_format_number(float(summary['Annual electricity saving']['value']))} MWh/year**."
        ),
        "",
        "## 1. Problem Definition",
        "",
        "The project objective is to use LNG cold energy to replace a conventional mechanical cooling load in a data center while satisfying the original assignment constraints on room conditions, chilled-water conditions, vaporizer duty, and pipeline distance.",
        "",
        "Key questions addressed in this study are:",
        "",
        "- What is the total cooling load of the target IDC?",
        "- What is the theoretical lower bound on power consumption?",
        "- How does a conventional R-134a reference cycle compare with an LNG-assisted system?",
        "- Which coolant is most suitable for the LNG-to-IDC secondary loop?",
        "- Is the system still feasible when the transport distance increases from 10 km to 35 km?",
        "- What are the annual energy, cost, and carbon implications?",
        "",
        "## 2. Design Basis",
        "",
        "- Assignment-derived design basis is recorded in [base.toml](../config/base.toml).",
        "- Full source registry is recorded in [sources.md](../docs/sources.md).",
        "- Engineering assumptions are recorded in [assumptions.md](../docs/assumptions.md).",
        "",
        "Core basis values:",
        "",
        f"- IDC total cooling load: **{_format_number(float(summary['IDC total cooling load']['value']))} kW**",
        f"- LNG / NG duty requirement: **{_format_number(float(summary['LNG vaporizer duty']['value']))} kW**",
        f"- Baseline R-134a compressor power: **{_format_number(float(summary['Baseline R-134a compressor power']['value']))} kW**",
        f"- Selected coolant: **{summary['Selected coolant']['value']}**",
        "",
        "## 3. Modeling Approach",
        "",
        "The codebase separates the project into load calculation, theoretical minimum power, baseline vapor-compression benchmarking, coolant screening, LNG vaporizer design, long-distance pipeline design, sensitivity studies, and annualized impact analysis.",
        "",
        "Main analysis outputs are generated by the following figures:",
        "",
        "![Load Breakdown](../output/figures/load_breakdown.png)",
        "",
        "![Baseline Cycle](../output/figures/baseline_cycle_ph.png)",
        "",
        "![Coolant Ranking](../output/figures/fluid_ranking.png)",
        "",
        "![HX Profile](../output/figures/hx_temperature_profile.png)",
        "",
        "![Pipeline Tradeoff](../output/figures/pipeline_tradeoff.png)",
        "",
        "## 4. Main Results",
        "",
        _markdown_table(
            pd.DataFrame(
                [
                    {"Metric": "Cooling load", "Value": f"{_format_number(float(summary['IDC total cooling load']['value']))} kW"},
                    {"Metric": "Theoretical minimum power", "Value": f"{_format_number(float(summary['Theoretical minimum power']['value']))} kW"},
                    {"Metric": "Baseline R-134a power", "Value": f"{_format_number(float(summary['Baseline R-134a compressor power']['value']))} kW"},
                    {"Metric": "Selected coolant", "Value": summary["Selected coolant"]["value"]},
                    {"Metric": "LNG loop pump power", "Value": f"{_format_number(float(summary['LNG system pump power']['value']))} kW"},
                    {"Metric": "Power saving", "Value": f"{_format_number(float(summary['Baseline-to-LNG power saving']['value']))} kW"},
                ]
            ),
            ["Metric", "Value"],
        ),
        "",
        "### 4.1 Coolant Selection",
        "",
        _markdown_table(
            alternatives.assign(
                screening_score=alternatives["screening_score"].map(lambda x: _format_number(x, 3)),
                pump_power_kw=alternatives["pump_power_kw"].map(_format_number),
                hx_shell_diameter_m=alternatives["hx_shell_diameter_m"].map(lambda x: _format_number(x, 3)),
                annual_cost_saving_krw=alternatives["annual_cost_saving_krw"].map(lambda x: _format_number(x / 1_000_000.0)),
            ),
            ["fluid", "screening_score", "pump_power_kw", "hx_shell_diameter_m", "annual_cost_saving_krw"],
        ),
        "",
        "The base-case optimum is **"
        f"{selected_alternative['fluid']}**, which minimizes loop pumping demand while preserving a feasible LNG vaporizer and pipeline design.",
        "",
        "### 4.2 Distance Sensitivity",
        "",
        _markdown_table(
            distance.assign(
                distance_km=distance["distance_km"].map(_format_number),
                pump_power_kw=distance["pump_power_kw"].map(_format_number),
                heat_gain_kw=distance["heat_gain_kw"].map(_format_number),
                thermal_margin_kw=distance["thermal_margin_kw"].map(_format_number),
            ),
            ["distance_km", "pump_power_kw", "heat_gain_kw", "thermal_margin_kw", "meets_idc_load"],
        ),
        "",
        "![Distance Sensitivity](../output/figures/pipeline_distance_sensitivity.png)",
        "",
        (
            f"The long-distance case at **{_format_number(long_distance['distance_km'])} km** is "
            f"**{'feasible' if bool(long_distance['meets_idc_load']) else 'not feasible'}** under the current duty margin. "
            f"The estimated maximum feasible one-way distance for the selected design is about "
            f"**{_format_number(distance['max_feasible_distance_m'].iloc[0] / 1000.0)} km**."
        ),
        "",
        "### 4.3 Supply Temperature Sensitivity",
        "",
        _markdown_table(
            temperature.assign(
                supply_temp_c=temperature["supply_temp_c"].map(_format_number),
                pump_power_kw=temperature["pump_power_kw"].map(_format_number),
                max_feasible_distance_km=temperature["max_feasible_distance_km"].map(_format_number),
            ),
            ["supply_temp_c", "selected_fluid", "pump_power_kw", "max_feasible_distance_km", "long_distance_meets_load", "status"],
        ),
        "",
        "![Supply Temperature Sensitivity](../output/figures/supply_temperature_sensitivity.png)",
        "",
        (
            f"The best pump-power point in the temperature sweep is **{_format_number(best_temp['supply_temp_c'])} C**, "
            f"which retains **{best_temp['selected_fluid']}** as the preferred fluid. "
            f"A warmer supply temperature of **-43.1 C** makes the 35 km case feasible, but only after switching to **R-600a** "
            f"and accepting a much larger pumping demand."
        ),
        "",
        "### 4.4 Annual Impact",
        "",
        _markdown_table(annual_report, ["metric", "value", "unit"]),
        "",
        "![Annual Impact](../output/figures/annual_impact_comparison.png)",
        "",
        "Allowable additional investment based on simple payback targets:",
        "",
        _markdown_table(
            payback.assign(
                allowable_incremental_capex_krw=payback["allowable_incremental_capex_krw"].map(lambda x: f"{_format_number(x / 1_000_000.0)} million KRW")
            ),
            ["payback_years", "allowable_incremental_capex_krw"],
        ),
        "",
        "## 5. Comparison with Legacy Excel Work",
        "",
        _markdown_table(
            legacy.assign(
                legacy_value_kw=legacy["legacy_value_kw"].map(_format_number),
                current_value_kw=legacy["current_value_kw"].map(_format_number),
                difference_kw=legacy["difference_kw"].map(_format_number),
                difference_percent=legacy["difference_percent"].map(lambda x: _format_number(x, 2)),
            ),
            ["metric", "legacy_value_kw", "current_value_kw", "difference_kw", "difference_percent"],
        ),
        "",
        "The recreated Python workflow is consistent with the historical spreadsheet at the system level, while making the assumptions, source traceability, and scenario studies explicit.",
        "",
        "## 6. Discussion",
        "",
        "- The base 10 km system is feasible with a large power reduction relative to the baseline compressor load.",
        "- The 35 km case is not feasible at the current duty margin for the base design point, which is itself a useful design conclusion rather than a failure.",
        "- A temperature-level shift can extend feasible distance, but it may require a different coolant and a higher loop pumping penalty.",
        "- The annualized economics are strong under the current boundary definition, but the comparison currently excludes detailed auxiliary equipment, maintenance, financing, and LNG infrastructure CAPEX.",
        "",
        "## 7. Limitations",
        "",
        "- LNG is modeled as pure methane in v1.",
        "- The IDC-side distribution network is simplified into a cooling-duty boundary rather than a full hydraulic network.",
        "- The economic comparison boundary is limited to baseline compressor power versus LNG loop pump power.",
        "- Detailed mechanical design checks such as stress, materials procurement, and controls integration are outside the present scope.",
        "",
        "## 8. Conclusion",
        "",
        (
            "The reconstructed project shows that an LNG cold-energy cooling concept can strongly reduce the modeled electrical demand of the reference data-center cooling system at the 10 km design point. "
            "The selected base design uses ammonia in the secondary loop, a shell-and-tube LNG vaporizer, and a 0.35 m supply/return pipeline pair. "
            "The design is energy-efficient, economically attractive under the present boundary, and transparent enough to support further report refinement and presentation work."
        ),
        "",
        "## References",
        "",
    ]
    for source_id in ["SRC-001", "SRC-004", "SRC-005", "SRC-006", "SRC-007", "SRC-008", "SRC-009", "SRC-010", "SRC-011", "SRC-012", "SRC-013", "SRC-014"]:
        if source_id in sources:
            entry = sources[source_id]
            report_lines.append(f"- **{source_id}** {entry['title']}. {entry['link_or_path']}")
    report_path.write_text("\n".join(report_lines), encoding="utf-8")
    return report_path


def build_presentation_script(project_root: Path) -> Path:
    output_dir = project_root / "output"
    deliverables_dir = ensure_directory(project_root / "deliverables")
    summary = _summary_map(output_dir)
    distance = pd.read_csv(output_dir / "distance_scenarios.csv")
    temperature = pd.read_csv(output_dir / "supply_temperature_sweep.csv")
    long_distance = distance.sort_values("distance_km").iloc[-1]
    high_temp = temperature[(temperature["status"] == "feasible") & (temperature["long_distance_meets_load"] == True)].head(1)

    script_lines = [
        "# Presentation Script",
        "",
        "## Slide 1. Title",
        "- State the project goal: replacing a conventional data-center cooling load with LNG cold energy.",
        "- Emphasize that the project was rebuilt as a reproducible code-based study.",
        "",
        "## Slide 2. Why This Problem Matters",
        "- Data centers concentrate electrical load and cooling demand.",
        "- LNG regasification discards large amounts of cold energy.",
        "- The project asks whether that cold energy can reduce cooling power demand.",
        "",
        "## Slide 3. Design Basis",
        f"- Total modeled cooling load is {_format_number(float(summary['IDC total cooling load']['value']))} kW.",
        "- Base transport distance is 10 km and the challenge case is 35 km.",
        "- The study keeps explicit source and assumption traceability.",
        "",
        "## Slide 4. Baseline Benchmark",
        f"- Theoretical minimum power is {_format_number(float(summary['Theoretical minimum power']['value']))} kW.",
        f"- Reference R-134a compressor power is {_format_number(float(summary['Baseline R-134a compressor power']['value']))} kW.",
        "- This baseline is the anchor for judging LNG-system benefit.",
        "",
        "## Slide 5. Proposed LNG Cooling Concept",
        "- LNG cold energy cools a secondary-loop refrigerant through a shell-and-tube vaporizer.",
        "- The secondary loop transports cooling duty from the terminal to the IDC.",
        "- The final base-case fluid is ammonia.",
        "",
        "## Slide 6. Coolant Screening Result",
        "- Among the feasible fluids, ammonia gives the lowest loop pumping power in the base case.",
        "- Propane and isobutane remain feasible alternatives but are materially weaker in pump power.",
        "- The selection is not arbitrary because the full ranking is reproduced by code.",
        "",
        "## Slide 7. Heat Exchanger Design",
        "- The LNG vaporizer is solved with a segmented enthalpy-based model.",
        "- The selected geometry is 500 tubes x 14 m with a shell diameter of about 0.723 m.",
        "- The minimum pinch is held at 10 K.",
        "",
        "## Slide 8. Pipeline Result",
        f"- Base-case LNG loop pumping power is {_format_number(float(summary['LNG system pump power']['value']))} kW.",
        f"- Estimated maximum feasible one-way distance is {_format_number(distance['max_feasible_distance_m'].iloc[0] / 1000.0)} km.",
        f"- Therefore the {int(long_distance['distance_km'])} km case is {'feasible' if bool(long_distance['meets_idc_load']) else 'not feasible'} at the base design point.",
        "",
        "## Slide 9. Sensitivity Insight",
        "- Distance sensitivity shows where the thermal margin collapses.",
        "- Supply-temperature sensitivity shows that 35 km can be recovered by moving to a warmer supply temperature.",
    ]
    if not high_temp.empty:
        row = high_temp.iloc[0]
        script_lines.append(
            f"- At {row['supply_temp_c']:.1f} C supply temperature, the design becomes feasible at 35 km with {row['selected_fluid']}."
        )
    script_lines.extend(
        [
            "",
            "## Slide 10. Annual Impact",
            f"- Annual electricity saving is {_format_number(float(summary['Annual electricity saving']['value']))} MWh/year.",
            f"- Annual electricity cost saving is {_format_number(float(summary['Annual electricity cost saving']['value']) / 1_000_000.0)} million KRW/year.",
            f"- Annual avoided indirect emissions are {_format_number(float(summary['Annual avoided indirect emissions']['value']))} tCO2/year.",
            "",
            "## Slide 11. Closing Message",
            "- The 10 km LNG cold-energy design is technically feasible and energetically attractive.",
            "- The 35 km case is the real design boundary test, not a simple extension of the base case.",
            "- The codebase now supports traceable report writing and future design refinement.",
        ]
    )
    script_path = deliverables_dir / "presentation_script.md"
    script_path.write_text("\n".join(script_lines), encoding="utf-8")
    return script_path


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str], footer: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(22)
    if footer:
        box = slide.shapes.add_textbox(Inches(0.4), Inches(6.7), Inches(12.0), Inches(0.35))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = footer
        paragraph.font.size = Pt(10)
        paragraph.alignment = PP_ALIGN.RIGHT


def _add_image_slide(prs: Presentation, title: str, bullets: list[str], image_path: Path, footer: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    textbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.9), Inches(4.5), Inches(5.7))
    frame = textbox.text_frame
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(20)
    slide.shapes.add_picture(str(image_path), Inches(4.9), Inches(1.0), width=Inches(8.1))
    if footer:
        box = slide.shapes.add_textbox(Inches(0.3), Inches(6.75), Inches(12.5), Inches(0.3))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = footer
        paragraph.font.size = Pt(10)
        paragraph.alignment = PP_ALIGN.RIGHT


def build_presentation(project_root: Path) -> Path:
    output_dir = project_root / "output"
    figure_dir = output_dir / "figures"
    deliverables_dir = ensure_directory(project_root / "deliverables")
    summary = _summary_map(output_dir)
    distance = pd.read_csv(output_dir / "distance_scenarios.csv")
    temperature = pd.read_csv(output_dir / "supply_temperature_sweep.csv")
    long_distance = distance.sort_values("distance_km").iloc[-1]
    best_temp = temperature[temperature["status"] == "feasible"].sort_values("pump_power_kw").iloc[0]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "LNG Cold-Energy Based IDC Cooling System",
        "Thermal System Design Project Reconstructed with Reproducible Python Models",
    )
    _add_bullet_slide(
        prs,
        "Why This Project Matters",
        [
            "Data centers require large and continuous cooling power.",
            "LNG regasification releases large amounts of usable cold energy.",
            "The project evaluates whether that cold energy can displace a conventional vapor-compression load.",
        ],
        footer="Sources: SRC-001, SRC-009, SRC-011",
    )
    _add_image_slide(
        prs,
        "Load and Baseline",
        [
            f"Modeled IDC cooling load: {_format_number(float(summary['IDC total cooling load']['value']))} kW",
            f"Theoretical minimum power: {_format_number(float(summary['Theoretical minimum power']['value']))} kW",
            f"Reference R-134a compressor power: {_format_number(float(summary['Baseline R-134a compressor power']['value']))} kW",
        ],
        figure_dir / "load_breakdown.png",
        footer="Sources: SRC-001, ASM-001 to ASM-011",
    )
    _add_image_slide(
        prs,
        "Baseline Benchmark Cycle",
        [
            "The baseline system is modeled as a simple R-134a vapor-compression cycle.",
            "This reference is used to quantify power savings from the LNG concept.",
            "The benchmark remains above the theoretical minimum, as expected.",
        ],
        figure_dir / "baseline_cycle_ph.png",
        footer="Sources: SRC-001, SRC-004, SRC-005",
    )
    _add_image_slide(
        prs,
        "Coolant Screening Result",
        [
            f"Selected coolant: {summary['Selected coolant']['value']}",
            "Ammonia gives the lowest loop pumping demand among feasible base-case options.",
            "Propane and isobutane remain viable fallback candidates.",
        ],
        figure_dir / "fluid_ranking.png",
        footer="Sources: SRC-003, SRC-008, ASM-017 to ASM-019",
    )
    _add_image_slide(
        prs,
        "LNG Vaporizer Design",
        [
            "Segmented enthalpy-based shell-and-tube design was used instead of a single constant-cp approximation.",
            "Selected geometry: 500 tubes x 14 m",
            "Minimum pinch maintained at 10 K",
        ],
        figure_dir / "hx_temperature_profile.png",
        footer="Sources: SRC-001, SRC-006, SRC-007",
    )
    _add_image_slide(
        prs,
        "Pipeline Design and Constraint",
        [
            f"Base-case loop pump power: {_format_number(float(summary['LNG system pump power']['value']))} kW",
            f"Estimated maximum feasible one-way distance: {_format_number(distance['max_feasible_distance_m'].iloc[0] / 1000.0)} km",
            f"The 35 km case is {'feasible' if bool(long_distance['meets_idc_load']) else 'not feasible'} at the base design point.",
        ],
        figure_dir / "pipeline_distance_sensitivity.png",
        footer="Sources: SRC-001, ASM-014 to ASM-016",
    )
    _add_image_slide(
        prs,
        "Temperature Sensitivity",
        [
            f"Best pump-power point: {_format_number(best_temp['supply_temp_c'])} C",
            f"Best fluid at that point: {best_temp['selected_fluid']}",
            "Warmer supply levels can extend feasible transport distance at the cost of larger pumping demand.",
        ],
        figure_dir / "supply_temperature_sensitivity.png",
        footer="Sources: ASM-028, ASM-029",
    )
    _add_image_slide(
        prs,
        "Annual Impact",
        [
            f"Annual electricity saving: {_format_number(float(summary['Annual electricity saving']['value']))} MWh/year",
            f"Annual cost saving: {_format_number(float(summary['Annual electricity cost saving']['value']) / 1_000_000.0)} million KRW/year",
            f"Annual avoided indirect emissions: {_format_number(float(summary['Annual avoided indirect emissions']['value']))} tCO2/year",
        ],
        figure_dir / "annual_impact_comparison.png",
        footer="Sources: SRC-013, SRC-014, ASM-030, ASM-032",
    )
    _add_image_slide(
        prs,
        "Comparison with Legacy Work",
        [
            "The recreated Python workflow reproduces the scale of the original spreadsheet benchmark.",
            "The new version adds traceability, validation, and full scenario sweeps.",
            "This turns a one-off project into a reusable design study.",
        ],
        figure_dir / "legacy_comparison.png",
        footer="Sources: SRC-010 and current Python workflow",
    )
    _add_bullet_slide(
        prs,
        "Conclusion",
        [
            "The LNG cold-energy concept is technically feasible at the 10 km base design point.",
            "The base design strongly reduces modeled electrical demand relative to the reference system.",
            "The 35 km case defines a real design boundary and motivates a trade-off between transport distance and operating temperature.",
        ],
        footer="Summary from reproducible project outputs",
    )

    pptx_path = deliverables_dir / "presentation_draft.pptx"
    prs.save(str(pptx_path))
    return pptx_path


def build_deliverables(project_root: Path) -> dict[str, Path]:
    report_path = build_report(project_root)
    script_path = build_presentation_script(project_root)
    presentation_path = build_presentation(project_root)
    return {
        "report": report_path,
        "script": script_path,
        "presentation": presentation_path,
    }
